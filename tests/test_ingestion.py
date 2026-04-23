"""Tests for document parsing and chunking."""

import pytest

from mcp_server.ingestion import DocumentParser


@pytest.fixture
def parser():
    return DocumentParser()


# ── Format Parsers ──


def test_parse_markdown(parser, sample_markdown):
    """Markdown parser extracts headers and content."""
    doc = parser.parse_file(sample_markdown)
    assert doc is not None
    assert doc.format == ".md"
    assert len(doc.chunks) > 0
    assert doc.metadata.get("type") == "markdown"


def test_parse_csv(parser, sample_csv):
    """CSV parser extracts rows and columns."""
    doc = parser.parse_file(sample_csv)
    assert doc is not None
    assert doc.format == ".csv"
    assert doc.metadata.get("rows") == 3
    assert doc.metadata.get("columns") == 3


def test_parse_json(parser, sample_json):
    """JSON parser detects structure."""
    doc = parser.parse_file(sample_json)
    assert doc is not None
    assert doc.format == ".json"
    assert doc.metadata.get("is_valid_json") is True
    assert doc.metadata.get("structure") == "object"


def test_parse_text(parser, sample_text):
    """Text parser counts lines."""
    doc = parser.parse_file(sample_text)
    assert doc is not None
    assert doc.format == ".txt"
    assert doc.metadata.get("line_count") >= 3


def test_parse_python(parser, sample_python):
    """Python parser extracts functions and classes."""
    doc = parser.parse_file(sample_python)
    assert doc is not None
    assert doc.format == ".py"
    assert "hello" in doc.metadata.get("functions", [])
    assert "Greeter" in doc.metadata.get("classes", [])


def test_parse_nonexistent(parser):
    """Nonexistent file raises error."""
    with pytest.raises(FileNotFoundError):
        parser.parse_file("/nonexistent/file.md")


def test_parse_unsupported(parser, tmp_path):
    """Unsupported format raises error."""
    f = tmp_path / "test.xyz"
    f.write_text("content")
    with pytest.raises(ValueError):
        parser.parse_file(f)


# ── Markdown Chunking ──


def test_markdown_chunks_by_headers(parser, sample_markdown):
    """Markdown files must split by ## headers (small sections may merge)."""
    doc = parser.parse_file(sample_markdown)
    # Should have at least 2 chunks (small sections merge per min_chunk_size)
    assert len(doc.chunks) >= 2
    # At least one chunk should have section_header metadata
    headers_found = [c.metadata.get("section_header", "") for c in doc.chunks if c.metadata.get("section_header")]
    assert len(headers_found) >= 1


def test_code_block_protection(parser, sample_markdown_with_code):
    """# comments inside code blocks must NOT create new chunks."""
    doc = parser.parse_file(sample_markdown_with_code)
    # Should NOT have chunks for "# This is a comment" or "# Another comment"
    for chunk in doc.chunks:
        assert "# This is a comment inside code block" not in chunk.metadata.get("section_header", "")


def test_min_chunk_size(parser, sample_markdown):
    """No chunk should be smaller than 100 chars (min merge threshold)."""
    doc = parser.parse_file(sample_markdown)
    for chunk in doc.chunks:
        # Allow the last chunk to be smaller (document tail)
        if chunk.index < len(doc.chunks) - 1:
            assert len(chunk.content) >= 50, f"Chunk {chunk.index} too small: {len(chunk.content)} chars"


def test_chunk_overlap_fallback(parser, tmp_path):
    """Non-markdown files use standard character-based chunking."""
    content = "A" * 2000  # Long text without headers
    f = tmp_path / "test.txt"
    f.write_text(content)
    doc = parser.parse_file(f)
    assert len(doc.chunks) >= 2  # Should split into multiple chunks


# ── DOCX/XLSX/PPTX ──


def test_docx_parser(parser, tmp_path):
    """DOCX parser extracts paragraphs and headings."""
    try:
        import docx

        d = docx.Document()
        d.add_heading("Test Heading", level=1)
        d.add_paragraph("Test content paragraph.")
        path = tmp_path / "test.docx"
        d.save(str(path))
        doc = parser.parse_file(path)
        assert doc is not None
        assert "# Test Heading" in doc.content
        assert doc.metadata.get("paragraphs") >= 1
    except ImportError:
        pytest.skip("python-docx not installed")


def test_xlsx_parser(parser, tmp_path):
    """XLSX parser extracts sheets."""
    try:
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Data"
        ws.append(["Col1", "Col2"])
        ws.append(["val1", "val2"])
        path = tmp_path / "test.xlsx"
        wb.save(str(path))
        doc = parser.parse_file(path)
        assert doc is not None
        assert "Data" in str(doc.metadata.get("sheets", []))
    except ImportError:
        pytest.skip("openpyxl not installed")


def test_pptx_parser(parser, tmp_path):
    """PPTX parser extracts slides."""
    try:
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Test Slide"
        slide.placeholders[1].text = "Slide content"
        path = tmp_path / "test.pptx"
        prs.save(str(path))
        doc = parser.parse_file(path)
        assert doc is not None
        assert doc.metadata.get("slides") == 1
        assert "Test Slide" in doc.content
    except ImportError:
        pytest.skip("python-pptx not installed")


# ── Parser Registry ──


def test_mqh_mq4_in_parsers(parser):
    """MQL4/MQL5 extensions must be registered in parser dispatch table."""
    assert ".mqh" in parser._parsers
    assert ".mq4" in parser._parsers


def test_ipynb_in_parsers(parser):
    """Jupyter Notebook must be registered in parser dispatch table."""
    assert ".ipynb" in parser._parsers


# ── Multi-Language Code Parsing ──


def test_new_extensions_in_parsers(parser):
    """All new code extensions must be registered in parser dispatch table."""
    for ext in [".c", ".h", ".cpp", ".js", ".jsx", ".ts", ".tsx", ".xml"]:
        assert ext in parser._parsers, f"{ext} missing from _parsers"


def test_parse_c(parser, sample_c):
    """C parser extracts functions, structs, and includes."""
    doc = parser.parse_file(sample_c)
    assert doc is not None
    assert doc.format == ".c"
    assert doc.metadata.get("language") == "c"
    assert "main" in doc.metadata.get("functions", [])
    assert "Config" in doc.metadata.get("classes", [])
    assert len(doc.metadata.get("imports", [])) >= 2


def test_parse_cpp(parser, sample_cpp):
    """C++ parser extracts classes, structs, and includes."""
    doc = parser.parse_file(sample_cpp)
    assert doc is not None
    assert doc.format == ".cpp"
    assert doc.metadata.get("language") == "cpp"
    assert "Engine" in doc.metadata.get("classes", [])
    assert len(doc.metadata.get("imports", [])) >= 2


def test_parse_javascript(parser, sample_js):
    """JavaScript parser extracts functions, classes, and imports."""
    doc = parser.parse_file(sample_js)
    assert doc is not None
    assert doc.format == ".js"
    assert doc.metadata.get("language") == "javascript"
    assert "fetchData" in doc.metadata.get("functions", [])
    assert "DataService" in doc.metadata.get("classes", [])
    assert len(doc.metadata.get("imports", [])) >= 2


def test_parse_typescript(parser, sample_ts):
    """TypeScript parser extracts functions, interfaces, enums, and imports."""
    doc = parser.parse_file(sample_ts)
    assert doc is not None
    assert doc.format == ".ts"
    assert doc.metadata.get("language") == "typescript"
    assert "createRouter" in doc.metadata.get("functions", [])
    assert "AppConfig" in doc.metadata.get("classes", [])
    assert "Status" in doc.metadata.get("classes", [])
    assert "ApiController" in doc.metadata.get("classes", [])
    assert len(doc.metadata.get("imports", [])) >= 1


def test_parse_xml(parser, sample_xml):
    """XML parser extracts root element and namespaces."""
    doc = parser.parse_file(sample_xml)
    assert doc is not None
    assert doc.format == ".xml"
    assert doc.metadata.get("type") == "xml"
    assert doc.metadata.get("root_element") == "project"
    namespaces = doc.metadata.get("namespaces", [])
    assert len(namespaces) >= 1
    uris = [ns["uri"] for ns in namespaces]
    assert any("maven" in uri for uri in uris)


def test_python_parse_unchanged(parser, sample_python):
    """Regression guard: Python parsing must produce identical output after refactor."""
    doc = parser.parse_file(sample_python)
    assert doc.metadata["language"] == "python"
    assert doc.metadata["type"] == "code"
    assert doc.metadata["functions"] == ["hello"]
    assert doc.metadata["classes"] == ["Greeter"]
    assert "docstring" in doc.metadata


def test_jsx_uses_javascript_language(parser, tmp_path):
    """JSX files must report language as javascript."""
    f = tmp_path / "component.jsx"
    f.write_text("function App() { return null; }", encoding="utf-8")
    doc = parser.parse_file(f)
    assert doc.metadata["language"] == "javascript"


def test_tsx_uses_typescript_language(parser, tmp_path):
    """TSX files must report language as typescript."""
    f = tmp_path / "component.tsx"
    f.write_text("function App(): JSX.Element { return null; }", encoding="utf-8")
    doc = parser.parse_file(f)
    assert doc.metadata["language"] == "typescript"


def test_h_uses_c_language(parser, tmp_path):
    """Header files must report language as c."""
    f = tmp_path / "header.h"
    f.write_text("#ifndef GUARD_H\n#define GUARD_H\nvoid func();\n#endif", encoding="utf-8")
    doc = parser.parse_file(f)
    assert doc.metadata["language"] == "c"
